"""
Presentation Styles for Nuvini Portfolio Reports

Reusable style constants, KPI box creators, and formatting helpers
for generating professional executive presentations.
"""

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from typing import Optional, Tuple


# =============================================================================
# Brand Colors
# =============================================================================

class NuviniColors:
    """Nuvini brand color palette"""
    # Primary
    PRIMARY_BLUE = RGBColor(0x1E, 0x3A, 0x8A)      # #1E3A8A - headers, titles
    DARK_BLUE = RGBColor(0x1B, 0x4F, 0x8C)         # #1B4F8C - darker variant
    MID_BLUE = RGBColor(0x3B, 0x82, 0xF6)          # #3B82F6 - charts, accents
    LIGHT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)        # #DBEAFE - KPI backgrounds

    # Status colors
    GREEN = RGBColor(0x10, 0xB9, 0x81)             # #10B981 - positive values
    RED = RGBColor(0xEF, 0x44, 0x44)               # #EF4444 - negative values
    ORANGE = RGBColor(0xF5, 0x9E, 0x0B)            # #F59E0B - caution/warning
    YELLOW = RGBColor(0xFB, 0xBF, 0x24)            # #FBBF24 - recovery

    # Neutrals
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DARK_TEXT = RGBColor(0x1F, 0x29, 0x37)         # #1F2937 - body text
    GRAY = RGBColor(0x6B, 0x72, 0x80)              # #6B7280 - secondary text
    LIGHT_GRAY = RGBColor(0xF8, 0xFA, 0xFC)        # #F8FAFC - backgrounds
    LIGHT_GREEN_BG = RGBColor(0xF0, 0xFD, 0xF4)    # #F0FDF4 - green tinted bg

    # Company-specific colors
    MERCOS = RGBColor(0x3B, 0x82, 0xF6)            # Blue
    EFFECTI = RGBColor(0x10, 0xB9, 0x81)           # Green
    IPE_DIGITAL = RGBColor(0xF5, 0x9E, 0x0B)       # Orange
    DATAHUB = RGBColor(0x63, 0x66, 0xF1)           # Indigo
    ONCLICK = RGBColor(0xEF, 0x44, 0x44)           # Red
    LEADLOVERS = RGBColor(0x6B, 0x72, 0x80)        # Gray


# Company color mapping
COMPANY_COLORS = {
    'Mercos': NuviniColors.MERCOS,
    'Effecti': NuviniColors.EFFECTI,
    'Ipê Digital': NuviniColors.IPE_DIGITAL,
    'Datahub': NuviniColors.DATAHUB,
    'OnClick': NuviniColors.ONCLICK,
    'Leadlovers': NuviniColors.LEADLOVERS,
}


# =============================================================================
# KPI Card Styles
# =============================================================================

def create_kpi_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    value: str,
    label: str,
    sublabel: str = None,
    accent_color: RGBColor = None,
    value_color: RGBColor = None
) -> None:
    """
    Create a professional KPI metric card with left accent border.

    Args:
        slide: The slide to add the card to
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        value: The main metric value (e.g., "R$115.3M")
        label: The metric label (e.g., "Gross Revenue (Jan-Jul)")
        sublabel: Optional subtitle (e.g., "All 6 Portfolio Companies")
        accent_color: Left border color (defaults to PRIMARY_BLUE)
        value_color: Value text color (defaults to PRIMARY_BLUE)
    """
    accent = accent_color or NuviniColors.PRIMARY_BLUE
    val_color = value_color or NuviniColors.PRIMARY_BLUE

    # Main background shape
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NuviniColors.LIGHT_BLUE
    shape.line.color.rgb = NuviniColors.LIGHT_BLUE
    shape.adjustments[0] = 0.05  # Subtle corner radius

    # Left accent border (thin rectangle on left side)
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top),
        Inches(0.06), Inches(height)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent
    accent_bar.line.fill.background()

    # Value text
    value_box = slide.shapes.add_textbox(
        Inches(left + 0.1), Inches(top + 0.1),
        Inches(width - 0.2), Inches(height * 0.5)
    )
    tf = value_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = value
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = val_color

    # Label text
    label_box = slide.shapes.add_textbox(
        Inches(left + 0.1), Inches(top + height * 0.5),
        Inches(width - 0.2), Inches(height * 0.25)
    )
    tf = label_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(10)
    run.font.color.rgb = NuviniColors.DARK_TEXT

    # Sublabel if provided
    if sublabel:
        sublabel_box = slide.shapes.add_textbox(
            Inches(left + 0.1), Inches(top + height * 0.75),
            Inches(width - 0.2), Inches(height * 0.2)
        )
        tf = sublabel_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = sublabel
        run.font.size = Pt(9)
        run.font.color.rgb = NuviniColors.GRAY


def create_kpi_row(
    slide,
    kpis: list,
    top: float = 1.0,
    start_left: float = 0.5,
    kpi_width: float = 2.4,
    kpi_height: float = 1.1,
    spacing: float = 0.12
) -> None:
    """
    Create a row of KPI cards evenly spaced.

    Args:
        slide: The slide to add cards to
        kpis: List of tuples (value, label, sublabel) or (value, label)
        top: Top position in inches
        start_left: Starting left position in inches
        kpi_width: Width of each KPI card
        kpi_height: Height of each KPI card
        spacing: Space between cards
    """
    for i, kpi in enumerate(kpis):
        left = start_left + i * (kpi_width + spacing)
        value = kpi[0]
        label = kpi[1]
        sublabel = kpi[2] if len(kpi) > 2 else None
        create_kpi_card(slide, left, top, kpi_width, kpi_height, value, label, sublabel)


# =============================================================================
# Section Boxes
# =============================================================================

def create_content_box(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    bullets: list,
    bg_color: RGBColor = None,
    border_color: RGBColor = None,
    title_color: RGBColor = None
) -> None:
    """
    Create a content box with title and bullet points.

    Args:
        slide: The slide to add the box to
        left, top, width, height: Position and size in inches
        title: Section title
        bullets: List of bullet point strings
        bg_color: Background color (defaults to LIGHT_GRAY)
        border_color: Border color (defaults to LIGHT_BLUE)
        title_color: Title color (defaults to PRIMARY_BLUE)
    """
    bg = bg_color or NuviniColors.LIGHT_GRAY
    border = border_color or NuviniColors.LIGHT_BLUE
    title_c = title_color or NuviniColors.PRIMARY_BLUE

    # Background shape
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = bg
    box.line.color.rgb = border
    box.adjustments[0] = 0.03

    # Content text box
    text_box = slide.shapes.add_textbox(
        Inches(left + 0.2), Inches(top + 0.15),
        Inches(width - 0.4), Inches(height - 0.3)
    )
    tf = text_box.text_frame
    tf.word_wrap = True

    # Title
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = title_c

    # Bullets
    for bullet in bullets:
        p = tf.add_paragraph()
        p.space_before = Pt(6)
        p.level = 0
        run = p.add_run()
        # Use Unicode bullet character
        run.text = f"\u2022  {bullet}"
        run.font.size = Pt(10)
        run.font.color.rgb = NuviniColors.DARK_TEXT


def create_insights_box(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    insights: list,
    icon_bullets: bool = True
) -> None:
    """
    Create a Key Insights box with icon-style bullets.
    Similar to create_content_box but with icons represented by symbols.
    """
    # Icons as Unicode characters
    icons = ['\u2713', '\u25B6', '\u2605', '\u25CF', '\u2714']  # checkmarks, arrows, stars, dots

    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = NuviniColors.LIGHT_GRAY
    box.line.color.rgb = NuviniColors.LIGHT_BLUE
    box.adjustments[0] = 0.03

    text_box = slide.shapes.add_textbox(
        Inches(left + 0.2), Inches(top + 0.15),
        Inches(width - 0.4), Inches(height - 0.3)
    )
    tf = text_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = NuviniColors.PRIMARY_BLUE

    for i, insight in enumerate(insights):
        p = tf.add_paragraph()
        p.space_before = Pt(5)

        if icon_bullets:
            # Icon character
            icon_run = p.add_run()
            icon_run.text = f"{icons[i % len(icons)]}  "
            icon_run.font.size = Pt(10)
            icon_run.font.color.rgb = NuviniColors.MID_BLUE

        # Insight text
        text_run = p.add_run()
        text_run.text = insight if icon_bullets else f"\u2022  {insight}"
        text_run.font.size = Pt(10)
        text_run.font.color.rgb = NuviniColors.DARK_TEXT


# =============================================================================
# Company Cards
# =============================================================================

def create_company_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    company_name: str,
    metrics: dict = None,
    bullets: list = None
) -> None:
    """
    Create a company performance card with colored border.

    Args:
        slide: Slide to add card to
        left, top, width, height: Position and size
        company_name: Company name
        metrics: Optional dict of metric name -> value
        bullets: Optional list of bullet points
    """
    company_color = COMPANY_COLORS.get(company_name, NuviniColors.GRAY)

    # Card background
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = NuviniColors.LIGHT_GRAY
    card.line.color.rgb = company_color
    card.line.width = Pt(2)
    card.adjustments[0] = 0.05

    # Company name header
    header_box = slide.shapes.add_textbox(
        Inches(left + 0.15), Inches(top + 0.1),
        Inches(width - 0.3), Inches(0.35)
    )
    tf = header_box.text_frame
    p = tf.paragraphs[0]

    # Colored bullet
    dot_run = p.add_run()
    dot_run.text = "\u25CF "
    dot_run.font.size = Pt(12)
    dot_run.font.color.rgb = company_color

    name_run = p.add_run()
    name_run.text = company_name
    name_run.font.size = Pt(12)
    name_run.font.bold = True
    name_run.font.color.rgb = NuviniColors.PRIMARY_BLUE

    # Metrics or bullets
    content_top = top + 0.4

    if metrics:
        # Two-column metrics layout
        metric_box = slide.shapes.add_textbox(
            Inches(left + 0.15), Inches(content_top),
            Inches(width - 0.3), Inches(height - 0.5)
        )
        tf = metric_box.text_frame
        tf.word_wrap = True

        for name, value in metrics.items():
            p = tf.add_paragraph()
            name_run = p.add_run()
            name_run.text = f"{name}: "
            name_run.font.size = Pt(10)
            name_run.font.color.rgb = NuviniColors.DARK_TEXT

            val_run = p.add_run()
            val_run.text = str(value)
            val_run.font.size = Pt(10)
            val_run.font.bold = True
            # Color based on value (green for positive, red for negative)
            if isinstance(value, str) and value.startswith('+'):
                val_run.font.color.rgb = NuviniColors.GREEN
            elif isinstance(value, str) and value.startswith('-'):
                val_run.font.color.rgb = NuviniColors.RED
            else:
                val_run.font.color.rgb = NuviniColors.PRIMARY_BLUE

    if bullets:
        bullet_box = slide.shapes.add_textbox(
            Inches(left + 0.15), Inches(content_top if not metrics else content_top + 0.4),
            Inches(width - 0.3), Inches(height - 0.6)
        )
        tf = bullet_box.text_frame
        tf.word_wrap = True

        for bullet in bullets:
            p = tf.add_paragraph()
            p.space_before = Pt(3)
            run = p.add_run()
            run.text = f"\u2022 {bullet}"
            run.font.size = Pt(9)
            run.font.color.rgb = NuviniColors.DARK_TEXT


def create_mom_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    company_name: str,
    revenue_change: float,
    ebitda_change: float
) -> None:
    """
    Create a month-over-month comparison card for a company.
    """
    company_color = COMPANY_COLORS.get(company_name, NuviniColors.GRAY)

    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = NuviniColors.LIGHT_GRAY
    card.line.color.rgb = company_color
    card.line.width = Pt(1.5)
    card.adjustments[0] = 0.05

    # Company name
    name_box = slide.shapes.add_textbox(
        Inches(left + 0.1), Inches(top + 0.08),
        Inches(width - 0.2), Inches(0.3)
    )
    tf = name_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = company_name
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = NuviniColors.PRIMARY_BLUE

    # Metrics
    metrics_box = slide.shapes.add_textbox(
        Inches(left + 0.1), Inches(top + 0.4),
        Inches(width - 0.2), Inches(0.5)
    )
    tf = metrics_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]

    # Revenue
    p.add_run().text = "Revenue: "
    rev_run = p.add_run()
    rev_run.text = f"{revenue_change:+.1f}%"
    rev_run.font.color.rgb = NuviniColors.GREEN if revenue_change >= 0 else NuviniColors.RED

    p.add_run().text = "    EBITDA: "
    ebitda_run = p.add_run()
    ebitda_run.text = f"{ebitda_change:+.1f}%"
    ebitda_run.font.color.rgb = NuviniColors.GREEN if ebitda_change >= 0 else NuviniColors.RED

    for run in p.runs:
        run.font.size = Pt(9)


# =============================================================================
# Headers and Footers
# =============================================================================

def create_slide_title(
    slide,
    title: str,
    subtitle: str = None,
    center: bool = True
) -> None:
    """
    Create a slide title with optional subtitle.
    """
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.25),
        Inches(12.33), Inches(0.6)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    if center:
        p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = NuviniColors.PRIMARY_BLUE

    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.85),
            Inches(12.33), Inches(0.4)
        )
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        if center:
            p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = subtitle
        run.font.size = Pt(16)
        run.font.color.rgb = NuviniColors.MID_BLUE


def add_slide_footer(slide, text: str) -> None:
    """
    Add a footer to the slide.
    """
    footer_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(7.1),
        Inches(12.33), Inches(0.3)
    )
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(9)
    run.font.color.rgb = NuviniColors.GRAY


# =============================================================================
# Status Helpers
# =============================================================================

def get_status(rule_of_40: float, yoy_growth: float, ebitda_margin: float) -> Tuple[str, RGBColor]:
    """
    Determine company status and color based on metrics.

    Returns:
        Tuple of (status_text, status_color)
    """
    if rule_of_40 >= 60:
        return "Strong", NuviniColors.GREEN
    elif rule_of_40 >= 40:
        if yoy_growth >= 0.40:
            return "Growth", NuviniColors.GREEN
        return "Strong", NuviniColors.GREEN
    elif yoy_growth < 0:
        if ebitda_margin > 0.35:
            return "Turnaround", NuviniColors.ORANGE
        return "Recovery", NuviniColors.YELLOW
    else:
        return "Developing", NuviniColors.MID_BLUE


def get_growth_color(value: float) -> RGBColor:
    """Get color for a growth value."""
    return NuviniColors.GREEN if value >= 0 else NuviniColors.RED
