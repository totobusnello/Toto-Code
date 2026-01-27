"""
Table Helpers for Nuvini Portfolio Presentations

Creates styled tables with margin bars, status badges, and proper formatting
for executive portfolio reports.
"""

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from typing import Dict, List, Tuple, Optional

from presentation_styles import NuviniColors, COMPANY_COLORS, get_status, get_growth_color


# =============================================================================
# Portfolio Summary Table
# =============================================================================

def create_portfolio_summary_table(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    company_data: List[Dict],
    portfolio_totals: Dict
) -> None:
    """
    Create the full portfolio summary table with status badges.

    Args:
        slide: Slide to add table to
        left, top, width, height: Position and size
        company_data: List of dicts with company metrics:
            - name: Company name
            - net_revenue: R$ value
            - ebitda: R$ value
            - ebitda_margin: Decimal (e.g., 0.388 for 38.8%)
            - yoy_growth: Decimal
            - gross_margin: Decimal
            - fcf: R$ value
            - rule_of_40: Float
            - status: String (Strong/Growth/Turnaround/Recovery)
        portfolio_totals: Dict with same keys for totals row
    """
    headers = [
        "Company", "Net Revenue\n(Jan-Jul)", "EBITDA\n(Jan-Jul)", "EBITDA\nMargin",
        "YoY\nGrowth", "Gross\nMargin", "FCF\n(Jan-Jul)", "Rule of\n40", "Status"
    ]

    rows = len(company_data) + 2  # header + companies + total
    cols = len(headers)

    table = slide.shapes.add_table(
        rows, cols,
        Inches(left), Inches(top),
        Inches(width), Inches(height / rows * rows)
    ).table

    # Column widths (proportional)
    col_widths = [1.3, 1.2, 1.2, 0.9, 0.9, 0.9, 1.1, 0.8, 0.9]
    total_width = sum(col_widths)
    for i, w in enumerate(col_widths):
        table.columns[i].width = Emu(int(Inches(width * w / total_width)))

    # Header row styling
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = NuviniColors.PRIMARY_BLUE

        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(9)
        para.font.bold = True
        para.font.color.rgb = NuviniColors.WHITE
        para.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Company rows
    for row_idx, company in enumerate(company_data, start=1):
        values = [
            company.get('name', ''),
            company.get('net_revenue', 'N/A'),
            company.get('ebitda', 'N/A'),
            company.get('ebitda_margin', 'N/A'),
            company.get('yoy_growth', 'N/A'),
            company.get('gross_margin', 'N/A'),
            company.get('fcf', 'N/A'),
            company.get('rule_of_40', 'N/A'),
            company.get('status', 'N/A'),
        ]

        for col_idx, value in enumerate(values):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)

            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(9)
            para.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            # Color coding
            if col_idx == 0:  # Company name
                para.font.bold = True
                para.font.color.rgb = NuviniColors.PRIMARY_BLUE
                para.alignment = PP_ALIGN.LEFT

            elif col_idx == 4:  # YoY Growth
                yoy = company.get('yoy_growth_raw', 0)
                para.font.color.rgb = get_growth_color(yoy)

            elif col_idx == 8:  # Status
                status = company.get('status', '')
                if status == 'Strong':
                    para.font.color.rgb = NuviniColors.GREEN
                elif status == 'Growth':
                    para.font.color.rgb = NuviniColors.GREEN
                elif status == 'Turnaround':
                    para.font.color.rgb = NuviniColors.ORANGE
                elif status == 'Recovery':
                    para.font.color.rgb = NuviniColors.YELLOW
                else:
                    para.font.color.rgb = NuviniColors.DARK_TEXT

        # Alternate row shading
        if row_idx % 2 == 0:
            for col_idx in range(cols):
                table.cell(row_idx, col_idx).fill.solid()
                table.cell(row_idx, col_idx).fill.fore_color.rgb = RGBColor(0xF9, 0xFA, 0xFB)

    # Total row
    total_row = rows - 1
    total_values = [
        "TOTAL PORTFOLIO",
        portfolio_totals.get('net_revenue', 'N/A'),
        portfolio_totals.get('ebitda', 'N/A'),
        portfolio_totals.get('ebitda_margin', 'N/A'),
        portfolio_totals.get('yoy_growth', 'N/A'),
        portfolio_totals.get('gross_margin', 'N/A'),
        portfolio_totals.get('fcf', 'N/A'),
        portfolio_totals.get('rule_of_40', 'N/A'),
        portfolio_totals.get('status', 'Strong'),
    ]

    for col_idx, value in enumerate(total_values):
        cell = table.cell(total_row, col_idx)
        cell.text = str(value)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NuviniColors.LIGHT_BLUE

        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(9)
        para.font.bold = True
        para.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE


# =============================================================================
# Profitability Table with Bars
# =============================================================================

def create_profitability_table(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    company_data: List[Dict]
) -> None:
    """
    Create a profitability table with visual margin bars.

    Args:
        slide: Slide to add table to
        left, top, width, height: Position and size
        company_data: List of dicts with:
            - name: Company name
            - gross_margin: Decimal (e.g., 0.726)
            - ebitda_margin: Decimal (e.g., 0.388)
    """
    headers = ["Company", "Gross Margin", "EBITDA Margin"]
    rows = len(company_data) + 1  # header + companies
    cols = len(headers)

    table = slide.shapes.add_table(
        rows, cols,
        Inches(left), Inches(top),
        Inches(width), Inches(height / rows * rows)
    ).table

    # Column widths
    table.columns[0].width = Emu(int(Inches(width * 0.25)))
    table.columns[1].width = Emu(int(Inches(width * 0.375)))
    table.columns[2].width = Emu(int(Inches(width * 0.375)))

    # Header row
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = NuviniColors.PRIMARY_BLUE

        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(10)
        para.font.bold = True
        para.font.color.rgb = NuviniColors.WHITE
        para.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Company rows with visual bars
    for row_idx, company in enumerate(company_data, start=1):
        name = company.get('name', '')
        gross_margin = company.get('gross_margin', 0)
        ebitda_margin = company.get('ebitda_margin', 0)

        # Company name
        cell = table.cell(row_idx, 0)
        cell.text = name
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(10)
        para.font.bold = True
        para.font.color.rgb = NuviniColors.PRIMARY_BLUE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Gross Margin with bar representation
        cell = table.cell(row_idx, 1)
        # We'll show the value and use unicode blocks to simulate a bar
        bar_width = int(gross_margin * 10)  # Scale to ~10 blocks max
        bar = '\u2588' * min(bar_width, 10)  # Full block character
        cell.text = f"{bar}  {gross_margin*100:.1f}%"

        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(10)
        para.font.color.rgb = NuviniColors.MID_BLUE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        # EBITDA Margin with bar
        cell = table.cell(row_idx, 2)
        bar_width = int(ebitda_margin * 10)
        bar = '\u2588' * min(bar_width, 10)
        cell.text = f"{bar}  {ebitda_margin*100:.1f}%"

        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(10)
        para.font.color.rgb = NuviniColors.GREEN
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Alternate row colors
        if row_idx % 2 == 0:
            for col in range(cols):
                table.cell(row_idx, col).fill.solid()
                table.cell(row_idx, col).fill.fore_color.rgb = RGBColor(0xF9, 0xFA, 0xFB)


# =============================================================================
# FCF Quarterly Table
# =============================================================================

def create_fcf_table(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    quarters: List[str],
    fcf_values: List[float],
    yoy_changes: List[Optional[float]]
) -> None:
    """
    Create quarterly FCF table with YoY changes.

    Args:
        slide: Slide to add table to
        left, top, width, height: Position and size
        quarters: Quarter labels (e.g., ["Q1 2024", "Q2 2024", ...])
        fcf_values: FCF values in millions
        yoy_changes: YoY change percentages (None if not applicable)
    """
    headers = ["Quarter", "FCF (R$M)", "YoY Change"]
    rows = len(quarters) + 1
    cols = len(headers)

    table = slide.shapes.add_table(
        rows, cols,
        Inches(left), Inches(top),
        Inches(width), Inches(height / rows * rows)
    ).table

    # Column widths
    for i in range(cols):
        table.columns[i].width = Emu(int(Inches(width / cols)))

    # Header row
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = NuviniColors.PRIMARY_BLUE

        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(10)
        para.font.bold = True
        para.font.color.rgb = NuviniColors.WHITE
        para.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    for row_idx, (quarter, fcf, yoy) in enumerate(zip(quarters, fcf_values, yoy_changes), start=1):
        # Quarter
        cell = table.cell(row_idx, 0)
        cell.text = quarter
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(10)
        para.font.bold = True
        para.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        # FCF Value
        cell = table.cell(row_idx, 1)
        cell.text = f"{fcf:.1f}"
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(10)
        para.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        # YoY Change
        cell = table.cell(row_idx, 2)
        if yoy is not None:
            cell.text = f"{yoy:+.1f}%"
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(10)
            para.font.color.rgb = get_growth_color(yoy)
        else:
            cell.text = "-"
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(10)
            para.font.color.rgb = NuviniColors.GRAY

        para.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Alternate row colors
        if row_idx % 2 == 0:
            for col in range(cols):
                table.cell(row_idx, col).fill.solid()
                table.cell(row_idx, col).fill.fore_color.rgb = RGBColor(0xF9, 0xFA, 0xFB)

        # Highlight current quarter
        if "2025" in quarter and "Q2" in quarter:
            for col in range(cols):
                table.cell(row_idx, col).fill.solid()
                table.cell(row_idx, col).fill.fore_color.rgb = NuviniColors.LIGHT_BLUE


# =============================================================================
# Company Rankings List
# =============================================================================

def create_rankings_list(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    rankings: List[Tuple[str, float, float, float]],
    title: str = "Company Rankings"
) -> None:
    """
    Create a company rankings list with Rule of 40 scores.

    Args:
        slide: Slide to add to
        left, top, width, height: Position and size
        rankings: List of tuples (company_name, rule_of_40_score, growth_pct, margin_pct)
        title: Section title
    """
    # Background box
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = NuviniColors.LIGHT_GRAY
    box.line.color.rgb = NuviniColors.LIGHT_BLUE
    box.adjustments[0] = 0.03

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(left + 0.15), Inches(top + 0.1),
        Inches(width - 0.3), Inches(0.35)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = NuviniColors.PRIMARY_BLUE

    # Rankings content
    content_box = slide.shapes.add_textbox(
        Inches(left + 0.15), Inches(top + 0.45),
        Inches(width - 0.3), Inches(height - 0.55)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, (company, score, growth, margin) in enumerate(rankings):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.space_before = Pt(6)

        # Company color bullet
        color = COMPANY_COLORS.get(company, NuviniColors.GRAY)
        bullet_run = p.add_run()
        bullet_run.text = "\u25CF "
        bullet_run.font.size = Pt(11)
        bullet_run.font.color.rgb = color

        # Company name
        name_run = p.add_run()
        name_run.text = company
        name_run.font.size = Pt(11)
        name_run.font.bold = True
        name_run.font.color.rgb = NuviniColors.DARK_TEXT

        # Score (right aligned via tabs/spaces)
        score_run = p.add_run()
        score_run.text = f"    {score:.1f}"
        score_run.font.size = Pt(11)
        score_run.font.bold = True
        score_run.font.color.rgb = NuviniColors.PRIMARY_BLUE

        # Details on next line
        p = tf.add_paragraph()
        p.space_before = Pt(1)
        detail_run = p.add_run()
        detail_run.text = f"   {growth:.1f}% growth + {margin:.1f}% margin"
        detail_run.font.size = Pt(9)
        detail_run.font.color.rgb = NuviniColors.GRAY
