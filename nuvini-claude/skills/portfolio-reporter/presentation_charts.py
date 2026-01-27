"""
Chart Helpers for Nuvini Portfolio Presentations

Creates waterfall charts, scatter plots, line charts, and bar charts
for executive portfolio reports.
"""

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData, XyChartData
from typing import Dict, List, Tuple, Optional

from presentation_styles import NuviniColors, COMPANY_COLORS


# =============================================================================
# Line Charts
# =============================================================================

def create_monthly_trend_chart(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    months: List[str],
    gross_revenue: List[float],
    net_revenue: List[float],
    ebitda: List[float],
    title: str = "Monthly Performance Trend"
) -> None:
    """
    Create a multi-series line chart showing monthly trends.

    Args:
        slide: Slide to add chart to
        left, top, width, height: Position and size in inches
        months: List of month labels (e.g., ["Jan", "Feb", "Mar"])
        gross_revenue: Gross revenue values in millions
        net_revenue: Net revenue values in millions
        ebitda: EBITDA values in millions
        title: Chart title
    """
    chart_data = CategoryChartData()
    chart_data.categories = months

    chart_data.add_series('Gross Revenue (R$ millions)', gross_revenue)
    chart_data.add_series('Net Revenue (R$ millions)', net_revenue)
    chart_data.add_series('EBITDA (R$ millions)', ebitda)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        Inches(left), Inches(top),
        Inches(width), Inches(height),
        chart_data
    ).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    # Style the series
    plot = chart.plots[0]
    plot.has_data_labels = False

    # Color each series
    colors = [NuviniColors.MID_BLUE, NuviniColors.GREEN, NuviniColors.ORANGE]
    for i, series in enumerate(chart.series):
        series.format.line.color.rgb = colors[i]
        series.format.line.width = Pt(2.5)
        series.smooth = True


def create_company_performance_chart(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    months: List[str],
    net_revenue: List[float],
    ebitda: List[float],
    company_name: str = None
) -> None:
    """
    Create a dual-series line chart for company deep dive.
    Shows Net Revenue and EBITDA trends.
    """
    chart_data = CategoryChartData()
    chart_data.categories = months

    chart_data.add_series('Monthly Net Revenue (R$ millions)', net_revenue)
    chart_data.add_series('Monthly EBITDA (R$ millions)', ebitda)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        Inches(left), Inches(top),
        Inches(width), Inches(height),
        chart_data
    ).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    plot = chart.plots[0]

    # Style series - solid line for revenue, dashed for EBITDA
    if len(chart.series) > 0:
        chart.series[0].format.line.color.rgb = NuviniColors.MID_BLUE
        chart.series[0].format.line.width = Pt(2.5)
        chart.series[0].smooth = True

    if len(chart.series) > 1:
        chart.series[1].format.line.color.rgb = NuviniColors.GREEN
        chart.series[1].format.line.width = Pt(2)
        chart.series[1].smooth = True


# =============================================================================
# Bar Charts
# =============================================================================

def create_horizontal_bar_chart(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    companies: List[str],
    values: List[float],
    title: str = "Net Revenue by Company",
    value_prefix: str = "R$",
    value_suffix: str = "M"
) -> None:
    """
    Create a horizontal bar chart with company-colored bars.

    Args:
        slide: Slide to add chart to
        left, top, width, height: Position and size
        companies: Company names (categories)
        values: Values in millions
        title: Chart title
    """
    chart_data = CategoryChartData()
    chart_data.categories = companies
    chart_data.add_series('Net Revenue (R$ millions)', values)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(left), Inches(top),
        Inches(width), Inches(height),
        chart_data
    ).chart

    chart.has_legend = False

    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.font.size = Pt(10)
    plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

    # Color bars by company
    series = chart.series[0]
    for i, company in enumerate(companies):
        point = series.points[i]
        color = COMPANY_COLORS.get(company, NuviniColors.MID_BLUE)
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = color


def create_quarterly_fcf_chart(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    quarters: List[str],
    fcf_values: List[float]
) -> None:
    """
    Create a vertical bar chart for quarterly FCF analysis.
    """
    chart_data = CategoryChartData()
    chart_data.categories = quarters
    chart_data.add_series('Free Cash Flow (R$ millions)', fcf_values)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(left), Inches(top),
        Inches(width), Inches(height),
        chart_data
    ).chart

    chart.has_legend = False

    plot = chart.plots[0]
    plot.has_data_labels = False

    # Color all bars blue
    series = chart.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = NuviniColors.MID_BLUE


# =============================================================================
# Scatter Charts
# =============================================================================

def create_rule_of_40_scatter(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    company_data: Dict[str, Tuple[float, float]],
    show_diagonal: bool = True
) -> None:
    """
    Create a Rule of 40 scatter plot.

    Args:
        slide: Slide to add chart to
        left, top, width, height: Position and size
        company_data: Dict mapping company name to (yoy_growth_pct, ebitda_margin_pct)
                     where values are percentages (e.g., 18.9 for 18.9%)
        show_diagonal: Whether to show the Rule of 40 diagonal line
    """
    chart_data = XyChartData()

    # Add each company as its own series for individual colors
    for company_name, (growth, margin) in company_data.items():
        series = chart_data.add_series(company_name)
        series.add_data_point(growth, margin)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.XY_SCATTER,
        Inches(left), Inches(top),
        Inches(width), Inches(height),
        chart_data
    ).chart

    chart.has_legend = False  # We'll add legend manually or skip it

    # Style each series with company color
    for i, (company_name, _) in enumerate(company_data.items()):
        if i < len(chart.series):
            series = chart.series[i]
            color = COMPANY_COLORS.get(company_name, NuviniColors.GRAY)
            series.marker.format.fill.solid()
            series.marker.format.fill.fore_color.rgb = color
            series.marker.size = 12

    # Set axis labels
    try:
        chart.category_axis.axis_title.text_frame.text = "YoY Growth Rate (%)"
        chart.value_axis.axis_title.text_frame.text = "EBITDA Margin (%)"
    except:
        pass  # Some versions of python-pptx don't support this well


def create_simple_scatter_placeholder(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str = "Rule of 40 Scatter Plot"
) -> None:
    """
    Create a placeholder for scatter chart with manual positioning.
    Use when XY charts are problematic.
    """
    from pptx.enum.shapes import MSO_SHAPE

    # Background box
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = NuviniColors.LIGHT_GRAY
    box.line.color.rgb = NuviniColors.LIGHT_BLUE
    box.adjustments[0] = 0.02

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(left + 0.2), Inches(top + 0.1),
        Inches(width - 0.4), Inches(0.4)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    from pptx.enum.text import PP_ALIGN
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = NuviniColors.PRIMARY_BLUE


# =============================================================================
# Waterfall Charts
# =============================================================================

def create_waterfall_chart(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    labels: List[str],
    values: List[float],
    is_subtotal: List[bool] = None,
    title: str = "Financial Waterfall Analysis"
) -> None:
    """
    Create a waterfall chart showing financial flow.

    Note: python-pptx doesn't have native waterfall support, so we simulate
    using a stacked bar chart with invisible base segments.

    Args:
        slide: Slide to add chart to
        left, top, width, height: Position and size
        labels: Category labels (e.g., ["Gross Revenue", "Deductions", "Net Revenue", ...])
        values: Values for each category (use negative for decreases)
        is_subtotal: Boolean list indicating which bars are subtotals (shown full from zero)
        title: Chart title
    """
    if is_subtotal is None:
        is_subtotal = [False] * len(labels)

    # Calculate base values for stacked effect
    bases = []
    current_base = 0

    for i, (val, subtotal) in enumerate(zip(values, is_subtotal)):
        if subtotal or i == 0:
            bases.append(0)
            current_base = val
        else:
            if val >= 0:
                bases.append(current_base)
                current_base += val
            else:
                current_base += val
                bases.append(current_base)

    chart_data = CategoryChartData()
    chart_data.categories = labels

    # Base series (invisible)
    chart_data.add_series('Base', bases)

    # Value series (visible)
    display_values = []
    for val, subtotal, base in zip(values, is_subtotal, bases):
        if subtotal:
            display_values.append(val)
        else:
            display_values.append(abs(val))

    chart_data.add_series('Values', display_values)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED,
        Inches(left), Inches(top),
        Inches(width), Inches(height),
        chart_data
    ).chart

    chart.has_legend = False

    # Make base series invisible
    base_series = chart.series[0]
    base_series.format.fill.background()
    base_series.format.line.fill.background()

    # Color the value series
    value_series = chart.series[1]
    for i, (val, subtotal) in enumerate(zip(values, is_subtotal)):
        point = value_series.points[i]
        point.format.fill.solid()

        if subtotal:
            point.format.fill.fore_color.rgb = NuviniColors.MID_BLUE
        elif val >= 0:
            point.format.fill.fore_color.rgb = NuviniColors.GREEN
        else:
            point.format.fill.fore_color.rgb = NuviniColors.RED


def create_simple_waterfall(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    data: Dict[str, Tuple[float, str]],
    title: str = "Financial Waterfall Analysis (Jan-Jul 2025)"
) -> None:
    """
    Simplified waterfall using standard column chart with visual styling.

    Args:
        data: Dict of label -> (value_in_millions, type: 'total'|'add'|'subtract')
    """
    labels = list(data.keys())
    values = [d[0] for d in data.values()]
    types = [d[1] for d in data.values()]

    chart_data = CategoryChartData()
    chart_data.categories = labels
    chart_data.add_series('R$ millions', values)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(left), Inches(top),
        Inches(width), Inches(height),
        chart_data
    ).chart

    chart.has_legend = False

    # Color by type
    series = chart.series[0]
    for i, type_ in enumerate(types):
        point = series.points[i]
        point.format.fill.solid()
        if type_ == 'total':
            point.format.fill.fore_color.rgb = NuviniColors.MID_BLUE
        elif type_ == 'add':
            point.format.fill.fore_color.rgb = NuviniColors.GREEN
        else:  # subtract
            point.format.fill.fore_color.rgb = NuviniColors.RED

    # Add title above chart
    from pptx.enum.text import PP_ALIGN
    title_box = slide.shapes.add_textbox(
        Inches(left), Inches(top - 0.35),
        Inches(width), Inches(0.35)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = NuviniColors.PRIMARY_BLUE
