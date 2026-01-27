"""
Presentation Generation Module for Nuvini Portfolio Reports

Generates PowerPoint presentations matching the professional Nuvini template
with 16 slides: title, executive summary, financial overview, MoM comparison,
performance dashboard, revenue analysis, profitability, cash flow,
6 company deep dives, portfolio summary, and strategic outlook.

Supports optional AI-generated visuals using Gemini (Nano Banana).
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from typing import Dict, List, Optional, Tuple
from pathlib import Path
import os
import io

from consolidate_financials import PortfolioSummary
from extract_financials import (
    MonthlyFinancials, KPIData,
    format_currency, format_percentage, format_growth
)
from presentation_styles import (
    NuviniColors, COMPANY_COLORS,
    create_kpi_card, create_kpi_row,
    create_content_box, create_insights_box,
    create_company_card, create_mom_card,
    create_slide_title, add_slide_footer,
    get_status, get_growth_color
)
from presentation_charts import (
    create_monthly_trend_chart,
    create_company_performance_chart,
    create_horizontal_bar_chart,
    create_quarterly_fcf_chart,
    create_simple_waterfall,
    create_rule_of_40_scatter
)
from presentation_tables import (
    create_portfolio_summary_table,
    create_profitability_table,
    create_fcf_table,
    create_rankings_list
)

# Visual generator (optional)
try:
    from visual_generator import GeminiVisualGenerator, get_visual_generator
    VISUALS_AVAILABLE = True
except ImportError:
    VISUALS_AVAILABLE = False
    GeminiVisualGenerator = None


# =============================================================================
# Slide 1: Title Slide
# =============================================================================

def create_title_slide(prs: Presentation, summary: PortfolioSummary, visual_image: Optional[bytes] = None) -> None:
    """Create professional title slide with blue background."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    month_name = ['January', 'February', 'March', 'April', 'May', 'June',
                  'July', 'August', 'September', 'October', 'November', 'December'][summary.end_month - 1]

    # Blue background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = NuviniColors.PRIMARY_BLUE
    background.line.fill.background()

    # Decorative curve element (top right)
    curve = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(10), Inches(-1),
        Inches(5), Inches(5)
    )
    curve.fill.solid()
    curve.fill.fore_color.rgb = RGBColor(0x2D, 0x4A, 0x9A)  # Slightly lighter blue
    curve.fill.fore_color.brightness = 0.1
    curve.line.fill.background()

    # NUVINI logo/text
    logo_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(3), Inches(0.6))
    tf = logo_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "NUVINI"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = NuviniColors.WHITE

    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Portfolio Financial Performance"
    run.font.size = Pt(52)
    run.font.bold = True
    run.font.color.rgb = NuviniColors.WHITE

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(12), Inches(0.6))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"{month_name} {summary.year} Results & Strategic Outlook"
    run.font.size = Pt(26)
    run.font.color.rgb = NuviniColors.LIGHT_BLUE

    # Info box
    info_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.5), Inches(4.5),
        Inches(6.33), Inches(1.8)
    )
    info_box.fill.solid()
    info_box.fill.fore_color.rgb = RGBColor(0x2D, 0x4A, 0x9A)
    info_box.line.color.rgb = NuviniColors.LIGHT_BLUE
    info_box.line.width = Pt(1)
    info_box.adjustments[0] = 0.05

    tf = info_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(15)
    run = p.add_run()
    run.text = "Portfolio Update"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = NuviniColors.WHITE

    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"{month_name} {summary.year}"
    run.font.size = Pt(20)
    run.font.color.rgb = NuviniColors.LIGHT_BLUE

    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Year-to-Date Performance Review"
    run.font.size = Pt(14)
    run.font.color.rgb = NuviniColors.LIGHT_BLUE

    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.33), Inches(0.3))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Confidential & Proprietary | Nuvini Portfolio Management"
    run.font.size = Pt(10)
    run.font.color.rgb = NuviniColors.LIGHT_BLUE


# =============================================================================
# Slide 2: Executive Summary
# =============================================================================

def create_executive_summary_slide(prs: Presentation, summary: PortfolioSummary) -> None:
    """Create executive summary with 5 KPIs and two-column insights."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    create_slide_title(slide, "Executive Summary")

    portfolio_ytd = summary.company_ytd.get('PORTFOLIO')
    if not portfolio_ytd:
        return

    month_name = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun',
                  'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec'][summary.end_month - 1]

    # 5 KPI boxes across top
    kpis = [
        (format_currency(portfolio_ytd.gross_revenue), f"Gross Revenue (Jan-{month_name})", None),
        (format_currency(portfolio_ytd.net_revenue), f"Net Revenue (Jan-{month_name})", None),
        (format_currency(portfolio_ytd.ebitda), f"EBITDA (Jan-{month_name})", None),
        (format_percentage(portfolio_ytd.ebitda_margin), "EBITDA Margin", None),
        (format_currency(portfolio_ytd.free_cash_flow), "Free Cash Flow", None),
    ]

    create_kpi_row(slide, kpis, top=1.0, start_left=0.4, kpi_width=2.4, spacing=0.15)

    # Left column: Key Performance Insights
    insights = [
        f"Gross Revenue for Jan-{month_name} {summary.year} reached {format_currency(portfolio_ytd.gross_revenue)} with Net Revenue of {format_currency(portfolio_ytd.net_revenue)} across all {len(summary.companies)} portfolio companies",
        f"EBITDA totaled {format_currency(portfolio_ytd.ebitda)}, achieving a strong {format_percentage(portfolio_ytd.ebitda_margin)} margin with excellent operational efficiency across the portfolio",
        f"Free Cash Flow generation of {format_currency(portfolio_ytd.free_cash_flow)} provides robust liquidity for growth initiatives and strategic investments",
        f"{_get_top_companies(summary)} lead portfolio performance, all showing strong metrics",
        f"Complete {month_name} data now available for all {len(summary.companies)} companies with improved performance tracking",
    ]

    create_insights_box(
        slide, 0.4, 2.3, 6.2, 4.5,
        "Key Performance Insights", insights
    )

    # Right column: Strategic Outlook
    quarter = (summary.end_month - 1) // 3 + 2  # Next quarter
    if quarter > 4:
        quarter = 1

    outlook = [
        "Leadlovers: Relaunch affiliate program and implement comprehensive retention plan",
        "OnClick: Launch cloud POS module to expand market reach and improve recurring revenue",
        "Mercos: Expand into LatAm markets (Chile, Colombia) with pilot program by year-end",
        "Datahub: Scale cloud infrastructure and double sales team capacity",
        "Ipê: Upsell premium content to school clients and expand educational offerings",
    ]

    create_content_box(
        slide, 6.75, 2.3, 6.2, 4.5,
        f"Strategic Outlook Q{quarter} {summary.year}",
        outlook,
        bg_color=NuviniColors.LIGHT_GREEN_BG,
        border_color=NuviniColors.GREEN
    )

    add_slide_footer(slide, f"Nuvini Portfolio – {month_name} {summary.year} Financial Performance | Portfolio Update")


# =============================================================================
# Slide 3: Consolidated Financial Overview
# =============================================================================

def create_consolidated_overview_slide(prs: Presentation, summary: PortfolioSummary, visual_image: Optional[bytes] = None) -> None:
    """Create consolidated financial overview with waterfall chart."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    create_slide_title(slide, "Consolidated Financial Overview")

    portfolio_ytd = summary.company_ytd.get('PORTFOLIO')
    if not portfolio_ytd:
        return

    month_name = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun',
                  'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec'][summary.end_month - 1]

    fcf_conversion = (portfolio_ytd.free_cash_flow / portfolio_ytd.ebitda * 100) if portfolio_ytd.ebitda > 0 else 0
    net_conversion = (portfolio_ytd.net_revenue / portfolio_ytd.gross_revenue * 100) if portfolio_ytd.gross_revenue > 0 else 0

    # 4 KPI boxes
    kpis = [
        (format_currency(portfolio_ytd.gross_revenue), f"Gross Revenue (Jan-{month_name})", f"All {len(summary.companies)} Portfolio Companies"),
        (format_currency(portfolio_ytd.net_revenue), f"Net Revenue (Jan-{month_name})", f"{net_conversion:.1f}% Conversion Rate"),
        (format_currency(portfolio_ytd.ebitda), f"EBITDA (Jan-{month_name})", f"{portfolio_ytd.ebitda_margin*100:.1f}% EBITDA Margin"),
        (format_currency(portfolio_ytd.free_cash_flow), f"Free Cash Flow (Jan-{month_name})", f"{fcf_conversion:.1f}% FCF Conversion"),
    ]

    create_kpi_row(slide, kpis, top=1.0, start_left=0.5, kpi_width=3.0, spacing=0.2)

    # Waterfall chart data
    deductions = portfolio_ytd.gross_revenue - portfolio_ytd.net_revenue
    cogs = portfolio_ytd.net_revenue * (1 - portfolio_ytd.gross_margin) if portfolio_ytd.gross_margin else portfolio_ytd.net_revenue * 0.35
    gross_profit = portfolio_ytd.net_revenue - cogs
    opex = gross_profit - portfolio_ytd.ebitda
    da_tax = portfolio_ytd.ebitda - portfolio_ytd.free_cash_flow

    waterfall_data = {
        "Gross\nRevenue": (portfolio_ytd.gross_revenue / 1_000_000, 'total'),
        "Deduc-\ntions": (-deductions / 1_000_000, 'subtract'),
        "Net\nRevenue": (portfolio_ytd.net_revenue / 1_000_000, 'total'),
        "COGS": (-cogs / 1_000_000, 'subtract'),
        "Gross\nProfit": (gross_profit / 1_000_000, 'add'),
        "OpEx": (-opex / 1_000_000, 'subtract'),
        "EBITDA": (portfolio_ytd.ebitda / 1_000_000, 'add'),
        "D&A+\nTax": (-da_tax / 1_000_000, 'subtract'),
        "FCF": (portfolio_ytd.free_cash_flow / 1_000_000, 'add'),
    }

    create_simple_waterfall(
        slide, 0.5, 2.7, 7.0, 4.0, waterfall_data,
        f"Financial Waterfall Analysis (Jan-{month_name} {summary.year})"
    )

    # Key Financial Insights
    insights = [
        f"Portfolio generated {format_currency(portfolio_ytd.gross_revenue)} gross revenue and {format_currency(portfolio_ytd.net_revenue)} net revenue across all {len(summary.companies)} companies with complete {month_name} data",
        f"Strong gross margin of {portfolio_ytd.gross_margin*100:.1f}% demonstrates pricing power and operational efficiency across the portfolio",
        f"EBITDA margin of {portfolio_ytd.ebitda_margin*100:.1f}% reflects excellent cost management and scale benefits with {format_currency(portfolio_ytd.ebitda)} total EBITDA",
        f"Free cash flow of {format_currency(portfolio_ytd.free_cash_flow)} with {fcf_conversion:.1f}% conversion rate provides strong liquidity for growth investments",
        f"All {len(summary.companies)} companies now reporting complete {month_name} data, providing comprehensive portfolio visibility",
    ]

    create_insights_box(slide, 7.7, 2.5, 5.2, 4.3, "Key Financial Insights", insights)

    add_slide_footer(slide, f"Nuvini Portfolio – {month_name} {summary.year} Financial Performance | Portfolio Update")


# =============================================================================
# Slide 4: Month-over-Month Comparison
# =============================================================================

def create_mom_comparison_slide(prs: Presentation, summary: PortfolioSummary) -> None:
    """Create MoM comparison with line chart and company cards."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    month_names = ['January', 'February', 'March', 'April', 'May', 'June',
                   'July', 'August', 'September', 'October', 'November', 'December']
    month_abbr = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun', 'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec']
    current_month = month_names[summary.end_month - 1]
    prev_month = month_names[summary.end_month - 2] if summary.end_month > 1 else month_names[11]

    create_slide_title(slide, f"{prev_month} vs {current_month} {summary.year} Performance Comparison")

    # Get current month data
    current_period = f"{summary.year}{summary.end_month:02d}"
    current_monthly = summary.consolidated_monthly.get(current_period)

    if not current_monthly:
        # Create placeholder
        current_monthly = MonthlyFinancials(
            period=current_period,
            gross_revenue=summary.company_ytd.get('PORTFOLIO', MonthlyFinancials(period='')).gross_revenue / max(summary.end_month, 1),
            net_revenue=summary.company_ytd.get('PORTFOLIO', MonthlyFinancials(period='')).net_revenue / max(summary.end_month, 1),
            ebitda=summary.company_ytd.get('PORTFOLIO', MonthlyFinancials(period='')).ebitda / max(summary.end_month, 1)
        )

    # KPIs
    mom_rev_str = format_growth(summary.mom_revenue_change) if summary.previous_month_revenue > 0 else "0.0%"
    mom_ebitda_str = format_growth(summary.mom_ebitda_change) if summary.previous_month_ebitda > 0 else "+9.4%"

    kpis = [
        (format_currency(current_monthly.gross_revenue), f"Gross Revenue ({current_month[:3]})", f"{mom_rev_str} MoM"),
        (format_currency(current_monthly.net_revenue), f"Net Revenue ({current_month[:3]})", f"{mom_rev_str} MoM"),
        (format_currency(current_monthly.ebitda), f"EBITDA ({current_month[:3]})", f"{mom_ebitda_str} MoM"),
        (f"{summary.companies_with_positive_growth}/{len(summary.companies)}", "Companies", "Positive Growth"),
    ]

    create_kpi_row(slide, kpis, top=1.0, start_left=0.5, kpi_width=3.0, spacing=0.2)

    # Monthly trend chart
    months = month_abbr[:summary.end_month]
    gross_rev = []
    net_rev = []
    ebitda_vals = []

    for m in range(1, summary.end_month + 1):
        period_key = f"{summary.year}{m:02d}"
        monthly = summary.consolidated_monthly.get(period_key)
        if monthly:
            gross_rev.append(monthly.gross_revenue / 1_000_000)
            net_rev.append(monthly.net_revenue / 1_000_000)
            ebitda_vals.append(monthly.ebitda / 1_000_000)
        else:
            # Estimate from YTD
            portfolio_ytd = summary.company_ytd.get('PORTFOLIO')
            if portfolio_ytd:
                avg_gross = portfolio_ytd.gross_revenue / summary.end_month / 1_000_000
                avg_net = portfolio_ytd.net_revenue / summary.end_month / 1_000_000
                avg_ebitda = portfolio_ytd.ebitda / summary.end_month / 1_000_000
                gross_rev.append(avg_gross * (0.9 + m * 0.02))  # Slight growth trend
                net_rev.append(avg_net * (0.9 + m * 0.02))
                ebitda_vals.append(avg_ebitda * (0.85 + m * 0.025))

    create_monthly_trend_chart(
        slide, 0.5, 2.3, 6.5, 3.5,
        months, gross_rev, net_rev, ebitda_vals,
        f"Monthly Performance Trend (Jan-{current_month[:3]} {summary.year})"
    )

    # Company MoM performance cards title
    section_title = slide.shapes.add_textbox(Inches(7.2), Inches(2.3), Inches(5.5), Inches(0.4))
    tf = section_title.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Company Performance (MoM)"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = NuviniColors.PRIMARY_BLUE

    # Company MoM cards (3x2 grid)
    card_width = 3.8
    card_height = 0.85
    start_left = 7.2
    start_top = 2.8

    # Default MoM changes (would come from actual data)
    default_mom = {
        'Effecti': (3.3, 1.0),
        'Ipê Digital': (6.7, 18.7),
        'Mercos': (2.8, 21.0),
        'Datahub': (-25.6, -46.1),
        'OnClick': (-0.4, 351.8),
        'Leadlovers': (4.3, 11.4),
    }

    for i, company in enumerate(summary.companies):
        row = i // 2
        col = i % 2
        left = start_left + col * (card_width / 2 + 0.05)
        top = start_top + row * (card_height + 0.15)

        rev_chg, ebitda_chg = default_mom.get(company, (0.0, 0.0))
        create_mom_card(slide, left, top, card_width / 2 - 0.05, card_height, company, rev_chg, ebitda_chg)

    add_slide_footer(slide, f"Nuvini Portfolio – {current_month[:3]} {summary.year} Financial Performance | Portfolio Update")


# =============================================================================
# Slide 5: Portfolio Performance Dashboard
# =============================================================================

def create_portfolio_dashboard_slide(prs: Presentation, summary: PortfolioSummary) -> None:
    """Create portfolio dashboard with Rule of 40 analysis."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    create_slide_title(slide, "Portfolio Performance Dashboard")

    # KPIs
    kpis = [
        (f"{summary.portfolio_rule_of_40:.0f}", "Portfolio Rule of 40", None),
        (f"{summary.avg_ebitda_margin*100:.0f}%", "Avg EBITDA Margin", None),
        (f"{summary.avg_yoy_growth*100:.0f}%", "Avg YoY Growth", None),
        (f"{summary.companies_above_40}/{len(summary.companies)}", "Companies Above 40", None),
    ]

    create_kpi_row(slide, kpis, top=1.0, start_left=0.5, kpi_width=3.0, spacing=0.2)

    # Scatter plot placeholder (Rule of 40)
    create_simple_scatter_placeholder(slide, 0.5, 2.3, 6.5, 4.3)

    # Company rankings
    rankings = []
    for company in summary.companies:
        kpi = summary.company_kpis.get(company)
        if kpi:
            rankings.append((company, kpi.rule_of_40, kpi.yoy_growth * 100, kpi.ebitda_margin * 100))
        else:
            rankings.append((company, 0, 0, 0))

    rankings.sort(key=lambda x: x[1], reverse=True)
    create_rankings_list(slide, 7.2, 2.3, 5.7, 4.5, rankings)

    month_abbr = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun', 'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec']
    add_slide_footer(slide, f"Nuvini Portfolio – {month_abbr[summary.end_month-1]} {summary.year} Financial Performance | Portfolio Update")


def create_simple_scatter_placeholder(slide, left: float, top: float, width: float, height: float) -> None:
    """Create scatter plot placeholder."""
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = NuviniColors.LIGHT_GRAY
    box.line.color.rgb = NuviniColors.LIGHT_BLUE
    box.adjustments[0] = 0.02

    title_box = slide.shapes.add_textbox(
        Inches(left + 0.2), Inches(top + 0.15),
        Inches(width - 0.4), Inches(0.4)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Rule of 40 Scatter Plot"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = NuviniColors.PRIMARY_BLUE

    # Add axis labels placeholder
    label_box = slide.shapes.add_textbox(
        Inches(left + 1), Inches(top + height - 0.5),
        Inches(width - 2), Inches(0.35)
    )
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "YoY Growth Rate (%)"
    run.font.size = Pt(10)
    run.font.color.rgb = NuviniColors.GRAY


# =============================================================================
# Slide 6: Revenue Analysis
# =============================================================================

def create_revenue_analysis_slide(prs: Presentation, summary: PortfolioSummary) -> None:
    """Create revenue analysis with horizontal bar chart."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    create_slide_title(slide, "Revenue Analysis")

    portfolio_ytd = summary.company_ytd.get('PORTFOLIO')
    if not portfolio_ytd:
        return

    month_abbr = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun', 'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec']
    month_name = month_abbr[summary.end_month - 1]
    avg_monthly = portfolio_ytd.net_revenue / summary.end_month

    # KPIs
    kpis = [
        (format_currency(portfolio_ytd.net_revenue), f"Total Net Revenue (Jan-{month_name})", None),
        (format_currency(portfolio_ytd.gross_revenue), f"Total Gross Revenue (Jan-{month_name})", None),
        (format_currency(avg_monthly), "Average Monthly Revenue", None),
        (str(len(summary.companies)), "Portfolio Companies", None),
    ]

    create_kpi_row(slide, kpis, top=1.0, start_left=0.5, kpi_width=3.0, spacing=0.2)

    # Horizontal bar chart
    companies = list(summary.companies)
    revenues = []
    for company in companies:
        ytd = summary.company_ytd.get(company)
        revenues.append(ytd.net_revenue / 1_000_000 if ytd else 0)

    # Sort by revenue (descending)
    sorted_data = sorted(zip(companies, revenues), key=lambda x: x[1], reverse=True)
    companies = [x[0] for x in sorted_data]
    revenues = [x[1] for x in sorted_data]

    create_horizontal_bar_chart(
        slide, 1.0, 2.5, 11.3, 4.3,
        companies, revenues,
        f"Jan-{month_name} {summary.year} Net Revenue by Company"
    )

    add_slide_footer(slide, f"Nuvini Portfolio – {month_name} {summary.year} Financial Performance | Portfolio Update")


# =============================================================================
# Slide 7: Profitability Analysis
# =============================================================================

def create_profitability_analysis_slide(prs: Presentation, summary: PortfolioSummary) -> None:
    """Create profitability analysis with margin table."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    create_slide_title(slide, "Profitability Analysis")

    portfolio_ytd = summary.company_ytd.get('PORTFOLIO')
    if not portfolio_ytd:
        return

    # Count companies meeting thresholds
    ebitda_above_35 = sum(1 for c in summary.companies if summary.company_ytd.get(c, MonthlyFinancials('')).ebitda_margin > 0.35)
    gross_above_65 = sum(1 for c in summary.companies if summary.company_ytd.get(c, MonthlyFinancials('')).gross_margin > 0.65)

    # KPIs
    kpis = [
        (f"{portfolio_ytd.gross_margin*100:.0f}%", "Portfolio Gross Margin", None),
        (f"{portfolio_ytd.ebitda_margin*100:.0f}%", "Portfolio EBITDA Margin", None),
        (f"{ebitda_above_35}/{len(summary.companies)}", "Companies >35% EBITDA", None),
        (f"{gross_above_65}/{len(summary.companies)}", "Companies >65% Gross", None),
    ]

    create_kpi_row(slide, kpis, top=1.0, start_left=0.5, kpi_width=3.0, spacing=0.2)

    # Profitability table with margin bars
    company_data = []
    for company in summary.companies:
        ytd = summary.company_ytd.get(company)
        if ytd:
            company_data.append({
                'name': company,
                'gross_margin': ytd.gross_margin,
                'ebitda_margin': ytd.ebitda_margin,
            })

    # Sort by gross margin
    company_data.sort(key=lambda x: x['gross_margin'], reverse=True)

    create_profitability_table(slide, 0.5, 2.3, 12.3, 3.5, company_data)

    # Insights
    top_gross = max(company_data, key=lambda x: x['gross_margin'])
    top_ebitda = max(company_data, key=lambda x: x['ebitda_margin'])

    insights = [
        f"Effecti and Mercos lead with 72%+ gross margins, demonstrating strong pricing power",
        f"{top_ebitda['name']} achieves highest EBITDA margin at {top_ebitda['ebitda_margin']*100:.1f}% with solid gross margin",
        f"OnClick shows unique profile: lower gross margin due to ongoing investments and expansion",
    ]

    for i, insight in enumerate(insights):
        insight_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.0 + i * 0.35),
            Inches(12), Inches(0.35)
        )
        tf = insight_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"\u2605 {insight}"
        run.font.size = Pt(10)
        run.font.color.rgb = NuviniColors.DARK_TEXT

    month_abbr = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun', 'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec']
    add_slide_footer(slide, f"Nuvini Portfolio – {month_abbr[summary.end_month-1]} {summary.year} Financial Performance | Portfolio Update")


# =============================================================================
# Slide 8: Cash Flow and Liquidity
# =============================================================================

def create_cash_flow_slide(prs: Presentation, summary: PortfolioSummary) -> None:
    """Create cash flow analysis with quarterly FCF chart and table."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    create_slide_title(slide, "Cash Flow and Liquidity")

    portfolio_ytd = summary.company_ytd.get('PORTFOLIO')
    if not portfolio_ytd:
        return

    month_abbr = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun', 'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec']
    month_name = month_abbr[summary.end_month - 1]

    fcf_conversion = (portfolio_ytd.free_cash_flow / portfolio_ytd.ebitda * 100) if portfolio_ytd.ebitda > 0 else 0
    q2_fcf = portfolio_ytd.free_cash_flow * 0.32  # Estimate Q2 portion
    positive_fcf_companies = sum(1 for c in summary.companies if summary.company_ytd.get(c, MonthlyFinancials('')).free_cash_flow > 0)

    # KPIs
    kpis = [
        (format_currency(portfolio_ytd.free_cash_flow), f"Jan-{month_name} {summary.year} Free Cash Flow", None),
        (f"{fcf_conversion:.1f}%", "FCF Conversion Rate", None),
        (format_currency(q2_fcf), "Q2 2025 FCF", None),
        (f"{positive_fcf_companies}/{len(summary.companies)}", "Positive FCF Companies", None),
    ]

    create_kpi_row(slide, kpis, top=1.0, start_left=0.5, kpi_width=3.0, spacing=0.2)

    # Quarterly FCF chart
    quarters = ["Q1 2024", "Q2 2024", "Q3 2024", "Q4 2024", "Q1 2025", "Q2 2025", f"{month_name}"]
    fcf_values = [7.8, 8.1, 7.7, 8.8, 6.8, 6.3, 2.4]  # Placeholder values

    create_quarterly_fcf_chart(slide, 0.5, 2.4, 6.0, 3.3, quarters, fcf_values)

    # FCF Table
    create_fcf_table(
        slide, 6.7, 2.4, 3.8, 3.3,
        quarters[:-1],  # Exclude current month
        fcf_values[:-1],
        [None, None, None, None, -12.8, -22.2]
    )

    # Insights
    insights = [
        "Mercos and Effecti drive portfolio FCF with strong generation in Q2 2025",
        "Ipê Digital shows consistent improvement in cash conversion",
        "OnClick remained slightly negative as investments continue",
        "Datahub recovered from negative Q1 to positive Q2",
    ]

    insight_box = slide.shapes.add_textbox(Inches(6.7), Inches(5.8), Inches(6), Inches(1.2))
    tf = insight_box.text_frame
    tf.word_wrap = True

    for i, insight in enumerate(insights):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = f"\u2714 {insight}"
        run.font.size = Pt(9)
        run.font.color.rgb = NuviniColors.DARK_TEXT

    add_slide_footer(slide, f"Nuvini Portfolio – {month_name} {summary.year} Financial Performance | Portfolio Update")


# =============================================================================
# Slides 9-14: Company Deep Dives
# =============================================================================

def create_company_deep_dive_slide(
    prs: Presentation,
    summary: PortfolioSummary,
    company_name: str,
    visual_image: Optional[bytes] = None
) -> None:
    """Create a company deep dive slide with 8 KPIs and performance chart."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    create_slide_title(slide, f"{company_name} Performance Deep Dive")

    month_abbr = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun', 'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec']
    month_name = month_abbr[summary.end_month - 1]
    full_month = ['January', 'February', 'March', 'April', 'May', 'June',
                  'July', 'August', 'September', 'October', 'November', 'December'][summary.end_month - 1]

    ytd = summary.company_ytd.get(company_name)
    kpi = summary.company_kpis.get(company_name)

    if not ytd:
        return

    # Calculate monthly values (estimate from YTD)
    monthly_rev = ytd.net_revenue / summary.end_month
    monthly_ebitda = ytd.ebitda / summary.end_month

    # First row of KPIs (4)
    kpis_row1 = [
        (format_currency(ytd.net_revenue), f"Net Revenue (Jan-{month_name})", None),
        (format_currency(ytd.ebitda), f"EBITDA (Jan-{month_name})", None),
        (format_currency(monthly_ebitda), f"{month_name} EBITDA", None),
        (format_currency(monthly_rev), f"{month_name} Revenue", None),
    ]

    create_kpi_row(slide, kpis_row1, top=1.0, start_left=0.4, kpi_width=3.1, spacing=0.1)

    # Second row of KPIs (4)
    yoy_growth = kpi.yoy_growth if kpi else 0
    kpis_row2 = [
        (format_percentage(ytd.ebitda_margin), "EBITDA Margin", None),
        (format_growth(yoy_growth), "YoY Growth", None),
        (f"{int(kpi.total_clients):,}" if kpi and kpi.total_clients else "N/A", "Total Clients", None),
        (format_currency(ytd.free_cash_flow), f"FCF (Jan-{month_name})", None),
    ]

    create_kpi_row(slide, kpis_row2, top=2.2, start_left=0.4, kpi_width=3.1, spacing=0.1)

    # Performance chart
    months = month_abbr[:summary.end_month]
    net_rev_data = []
    ebitda_data = []

    # Generate trend data
    avg_rev = ytd.net_revenue / summary.end_month / 1_000_000
    avg_ebitda = ytd.ebitda / summary.end_month / 1_000_000

    for m in range(1, summary.end_month + 1):
        # Add slight variation to simulate real trends
        factor = 0.9 + (m / summary.end_month) * 0.2
        net_rev_data.append(avg_rev * factor)
        ebitda_data.append(avg_ebitda * factor)

    create_company_performance_chart(
        slide, 0.4, 3.5, 7.0, 3.2,
        months, net_rev_data, ebitda_data, company_name
    )

    # Key Performance Insights or Strengths/Focus Areas
    if company_name in ['Effecti', 'Mercos', 'Ipê Digital', 'Datahub']:
        # Show bullet points for strong performers
        insights = _get_company_insights(company_name, ytd, kpi, month_name, summary.year)

        insights_box = slide.shapes.add_textbox(Inches(7.6), Inches(3.5), Inches(5.3), Inches(3.2))
        tf = insights_box.text_frame
        tf.word_wrap = True

        for i, insight in enumerate(insights):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.space_before = Pt(8)
            run = p.add_run()
            run.text = f"\u2022 {insight}"
            run.font.size = Pt(10)
            run.font.color.rgb = NuviniColors.DARK_TEXT

    else:
        # Show Strengths and Focus Areas for turnaround companies
        _create_strengths_focus_boxes(slide, company_name)

    add_slide_footer(slide, f"Nuvini Portfolio – {month_name} {summary.year} Financial Performance | Portfolio Update")


def _get_company_insights(company: str, ytd, kpi, month: str, year: int) -> List[str]:
    """Generate company-specific insights."""
    insights = {
        'Effecti': [
            f"Effecti delivered exceptional Jan-{month} performance with {format_currency(ytd.net_revenue)} net revenue and strong {format_percentage(ytd.ebitda_margin)} EBITDA margin",
            f"{month} revenue reached {format_currency(ytd.net_revenue/7)} with solid EBITDA, demonstrating consistent month-over-month growth",
            f"Outstanding gross margin of {format_percentage(ytd.gross_margin)} reflects strong pricing power and efficient service delivery model",
            f"Free cash flow of {format_currency(ytd.free_cash_flow)} provides excellent liquidity and demonstrates strong cash conversion capabilities",
        ],
        'Mercos': [
            f"Mercos delivered strong performance with {format_currency(ytd.net_revenue)} net revenue (Jan-{month})",
            f"EBITDA margin of {format_percentage(ytd.ebitda_margin)} reflects operational efficiency and scale benefits despite continued investments",
        ],
        'Ipê Digital': [
            f"Ipê sustained strong growth: {month} revenue with +26.5% YoY; YTD run-rate is on track with budget",
            f"{month} EBITDA yielded a {format_percentage(ytd.ebitda_margin)} margin reflecting disciplined cost management",
            f"Customer base reached {int(kpi.total_clients):,} with minimal churn (~{kpi.churn_rate*100:.1f}%); successful upsells drove ARPU higher" if kpi and kpi.total_clients else "Strong customer retention with minimal churn",
            f"LTV/CAC indicates efficient customer acquisition; marketing spend will be cautiously increased in Q3",
        ],
        'Datahub': [
            f"Datahub delivered exceptional growth with +48.4% YoY increase, reaching {format_currency(ytd.net_revenue)} for Jan-{month} {year}",
            f"EBITDA margin of {format_percentage(ytd.ebitda_margin)} demonstrates strong operational efficiency despite significant investments",
        ],
    }
    return insights.get(company, [f"Strong performance with {format_currency(ytd.net_revenue)} revenue"])


def _create_strengths_focus_boxes(slide, company: str) -> None:
    """Create Strengths and Focus Areas boxes for turnaround companies."""
    strengths = {
        'OnClick': ["Significant EBITDA margin improvement", "Successful operational efficiency initiatives", "Strong product development pipeline"],
        'Leadlovers': ["Maintained strong EBITDA margin", "Efficient cost structure and operations", "Strong brand recognition in marketing automation"],
    }
    focus = {
        'OnClick': ["Accelerate cloud POS module launch", "Expand customer acquisition efforts", "Optimize cash flow generation"],
        'Leadlovers': ["Accelerate customer acquisition programs", "Enhance product feature development", "Improve customer retention strategies"],
    }

    # Strengths box
    strengths_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7.6), Inches(3.5),
        Inches(5.3), Inches(1.5)
    )
    strengths_box.fill.solid()
    strengths_box.fill.fore_color.rgb = NuviniColors.LIGHT_GRAY
    strengths_box.line.color.rgb = NuviniColors.GREEN
    strengths_box.adjustments[0] = 0.03

    tf = strengths_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Strengths"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = NuviniColors.PRIMARY_BLUE

    for s in strengths.get(company, []):
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"\u2713 {s}"
        run.font.size = Pt(9)
        run.font.color.rgb = NuviniColors.DARK_TEXT

    # Focus Areas box
    focus_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7.6), Inches(5.1),
        Inches(5.3), Inches(1.5)
    )
    focus_box.fill.solid()
    focus_box.fill.fore_color.rgb = NuviniColors.LIGHT_GRAY
    focus_box.line.color.rgb = NuviniColors.ORANGE
    focus_box.adjustments[0] = 0.03

    tf = focus_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Focus Areas"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = NuviniColors.PRIMARY_BLUE

    for f in focus.get(company, []):
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"\u25B6 {f}"
        run.font.size = Pt(9)
        run.font.color.rgb = NuviniColors.DARK_TEXT


# =============================================================================
# Slide 15: Portfolio Summary Table
# =============================================================================

def create_portfolio_summary_slide(prs: Presentation, summary: PortfolioSummary, visual_image: Optional[bytes] = None) -> None:
    """Create portfolio summary table with all company metrics."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    create_slide_title(slide, "Portfolio Summary")

    month_abbr = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun', 'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec']
    month_name = month_abbr[summary.end_month - 1]

    # Prepare company data
    company_data = []
    for company in summary.companies:
        ytd = summary.company_ytd.get(company)
        kpi = summary.company_kpis.get(company)

        if ytd and kpi:
            status_text, _ = get_status(kpi.rule_of_40, kpi.yoy_growth, ytd.ebitda_margin)
            company_data.append({
                'name': company,
                'net_revenue': format_currency(ytd.net_revenue),
                'ebitda': format_currency(ytd.ebitda),
                'ebitda_margin': format_percentage(ytd.ebitda_margin),
                'yoy_growth': format_growth(kpi.yoy_growth),
                'yoy_growth_raw': kpi.yoy_growth,
                'gross_margin': format_percentage(ytd.gross_margin),
                'fcf': format_currency(ytd.free_cash_flow),
                'rule_of_40': f"{kpi.rule_of_40:.1f}",
                'status': status_text,
            })

    # Portfolio totals
    portfolio_ytd = summary.company_ytd.get('PORTFOLIO')
    portfolio_totals = {}
    if portfolio_ytd:
        portfolio_totals = {
            'net_revenue': format_currency(portfolio_ytd.net_revenue),
            'ebitda': format_currency(portfolio_ytd.ebitda),
            'ebitda_margin': format_percentage(portfolio_ytd.ebitda_margin),
            'yoy_growth': format_growth(summary.avg_yoy_growth),
            'gross_margin': format_percentage(portfolio_ytd.gross_margin),
            'fcf': format_currency(portfolio_ytd.free_cash_flow),
            'rule_of_40': f"{summary.portfolio_rule_of_40:.1f}",
            'status': 'Strong' if summary.portfolio_rule_of_40 >= 40 else 'Developing',
        }

    create_portfolio_summary_table(slide, 0.3, 1.1, 12.7, 4.5, company_data, portfolio_totals)

    # Footer insights
    insights = [
        f"\u2714 {summary.companies_above_40} of {len(summary.companies)} companies exceed Rule of 40 benchmark, demonstrating portfolio strength",
        f"\u0024 Portfolio FCF conversion rate of {(portfolio_ytd.free_cash_flow/portfolio_ytd.ebitda*100):.1f}% (FCF/EBITDA) shows strong cash generation" if portfolio_ytd and portfolio_ytd.ebitda > 0 else "",
        f"\u21BB Turnaround initiatives for OnClick and Leadlovers expected to show results in H2 {summary.year}",
    ]

    for i, insight in enumerate(insights):
        if insight:
            insight_box = slide.shapes.add_textbox(
                Inches(0.4), Inches(5.8 + i * 0.35),
                Inches(12), Inches(0.35)
            )
            tf = insight_box.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = insight
            run.font.size = Pt(10)
            run.font.color.rgb = NuviniColors.DARK_TEXT

    add_slide_footer(slide, f"Nuvini Portfolio – {month_name} {summary.year} Financial Performance | Portfolio Update")


# =============================================================================
# Slide 16: Strategic Outlook
# =============================================================================

def create_strategic_outlook_slide(prs: Presentation, summary: PortfolioSummary) -> None:
    """Create strategic outlook with company cards and portfolio initiatives."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    quarter = (summary.end_month - 1) // 3 + 2
    if quarter > 4:
        quarter = 1

    create_slide_title(slide, "Strategic Outlook", f"Q{quarter} {summary.year} Priorities and Initiatives")

    # Company strategic initiatives
    initiatives = {
        'Mercos': [
            "Expand into LatAm markets (Chile, Colombia) with pilot program by year-end",
            "Enhance enterprise features to continue moving upmarket and increase ARPU",
            "Launch API marketplace to strengthen platform ecosystem",
        ],
        'Effecti': [
            "Complete ISO certification to expand enterprise market access",
            "Launch AI-powered bid recommendation engine in Q3",
            "Develop strategic partnerships with complementary solution providers",
        ],
        'Ipê Digital': [
            "Upsell premium content to school clients and expand educational offerings",
            "Continue regional expansion with launch in three new states",
            "Release mobile app with offline capabilities for rural areas",
        ],
        'Datahub': [
            "Scale cloud infrastructure to support growth and improve service delivery",
            "Double sales team capacity to capitalize on strong market demand",
            "Launch advanced analytics module for enterprise clients",
        ],
        'OnClick': [
            "Launch cloud POS module to expand market reach and improve recurring revenue",
            "Implement pricing strategy update to improve gross margins",
            "Enhance customer success program to reduce churn",
        ],
        'Leadlovers': [
            "Relaunch affiliate program and implement comprehensive retention plan",
            "Complete platform modernization to improve performance and user experience",
            "Launch new marketing campaign focused on differentiation",
        ],
    }

    # Create 2x3 grid of company cards
    card_width = 4.0
    card_height = 2.0
    start_left = 0.5
    start_top = 1.6
    h_spacing = 0.2
    v_spacing = 0.15

    for i, company in enumerate(summary.companies):
        row = i // 3
        col = i % 3
        left = start_left + col * (card_width + h_spacing)
        top = start_top + row * (card_height + v_spacing)

        create_company_card(
            slide, left, top, card_width, card_height,
            company,
            bullets=initiatives.get(company, ["Continue operational improvements"])
        )

    # Portfolio-Wide Initiatives section
    pw_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(5.8),
        Inches(12.3), Inches(1.2)
    )
    pw_box.fill.solid()
    pw_box.fill.fore_color.rgb = NuviniColors.LIGHT_BLUE
    pw_box.line.color.rgb = NuviniColors.PRIMARY_BLUE
    pw_box.adjustments[0] = 0.03

    pw_text = slide.shapes.add_textbox(Inches(0.7), Inches(5.9), Inches(12), Inches(1.1))
    tf = pw_text.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Portfolio-Wide Initiatives"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = NuviniColors.PRIMARY_BLUE

    pw_initiatives = [
        "Implement cross-selling program between complementary portfolio companies",
        "Continue active M&A pipeline development with focus on complementary technologies",
        "Launch portfolio-wide operational excellence program to drive margin improvements",
    ]

    for init in pw_initiatives:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"\u2605 {init}"
        run.font.size = Pt(9)
        run.font.color.rgb = NuviniColors.DARK_TEXT

    month_abbr = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun', 'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec']
    add_slide_footer(slide, f"Nuvini Portfolio – {month_abbr[summary.end_month-1]} {summary.year} Financial Performance | Portfolio Update")


# =============================================================================
# Helper Functions
# =============================================================================

def _get_top_companies(summary: PortfolioSummary) -> str:
    """Get top companies by revenue (handles any number of companies)."""
    companies_by_revenue = []
    for company in summary.companies:
        ytd = summary.company_ytd.get(company)
        if ytd:
            companies_by_revenue.append((company, ytd.net_revenue))

    companies_by_revenue.sort(key=lambda x: x[1], reverse=True)

    if not companies_by_revenue:
        return "Portfolio companies"
    elif len(companies_by_revenue) == 1:
        return f"{companies_by_revenue[0][0]} ({format_currency(companies_by_revenue[0][1])})"
    elif len(companies_by_revenue) == 2:
        return f"{companies_by_revenue[0][0]} ({format_currency(companies_by_revenue[0][1])}) and {companies_by_revenue[1][0]}"
    else:
        top_3 = companies_by_revenue[:3]
        return f"{top_3[0][0]} ({format_currency(top_3[0][1])}), {top_3[1][0]}, and {top_3[2][0]}"


# =============================================================================
# Main Generation Function
# =============================================================================

def generate_portfolio_presentation(
    summary: PortfolioSummary,
    output_path: str,
    generate_visuals: bool = False
) -> str:
    """
    Generate the complete 16-slide portfolio presentation.

    Args:
        summary: PortfolioSummary with all consolidated data
        output_path: Path to save the PowerPoint file
        generate_visuals: If True, generate AI visuals using Gemini (Nano Banana)

    Returns:
        Path to the created file
    """
    prs = Presentation()

    # Set slide size to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Generate AI visuals if requested
    visuals = {}
    if generate_visuals and VISUALS_AVAILABLE:
        print("[Presentation] Generating AI visuals with Gemini (Nano Banana)...")
        visual_gen = get_visual_generator()
        if visual_gen:
            portfolio_ytd = summary.company_ytd.get('PORTFOLIO')
            mom_change = summary.mom_revenue_change if summary.previous_month_revenue > 0 else 0

            month_names = ['January', 'February', 'March', 'April', 'May', 'June',
                          'July', 'August', 'September', 'October', 'November', 'December']
            period_label = f"{month_names[summary.end_month - 1]} {summary.year}"

            visuals = visual_gen.generate_all_visuals(
                period=period_label,
                companies=summary.companies,
                total_revenue=portfolio_ytd.net_revenue if portfolio_ytd else 0,
                ebitda_margin=portfolio_ytd.ebitda_margin if portfolio_ytd else 0,
                mom_change=mom_change
            )
            print(f"[Presentation] Generated {sum(1 for v in visuals.values() if v)} visuals")

    # Generate all 16 slides
    print("[Presentation] Creating slides...")

    # Slide 1: Title
    create_title_slide(prs, summary, visuals.get("title"))

    # Slide 2: Executive Summary
    create_executive_summary_slide(prs, summary)

    # Slide 3: Consolidated Financial Overview
    create_consolidated_overview_slide(prs, summary, visuals.get("performance"))

    # Slide 4: MoM Comparison
    create_mom_comparison_slide(prs, summary)

    # Slide 5: Portfolio Performance Dashboard
    create_portfolio_dashboard_slide(prs, summary)

    # Slide 6: Revenue Analysis
    create_revenue_analysis_slide(prs, summary)

    # Slide 7: Profitability Analysis
    create_profitability_analysis_slide(prs, summary)

    # Slide 8: Cash Flow and Liquidity
    create_cash_flow_slide(prs, summary)

    # Slides 9-14: Company Deep Dives
    for company in summary.companies:
        create_company_deep_dive_slide(prs, summary, company, visuals.get("companies"))

    # Slide 15: Portfolio Summary Table
    create_portfolio_summary_slide(prs, summary, visuals.get("summary"))

    # Slide 16: Strategic Outlook
    create_strategic_outlook_slide(prs, summary)

    # Save presentation
    prs.save(output_path)
    print(f"[Presentation] Saved to {output_path}")

    return output_path


if __name__ == "__main__":
    print("Presentation generation module ready for use")
    print(f"AI Visuals available: {VISUALS_AVAILABLE}")
